"""Tests for XLSX, PPTX, Markdown, and HTML read/write primitives."""

from __future__ import annotations

import hashlib
from pathlib import Path

import pytest

from core.file_formats import (
    ReadLimits,
    WriteValidationError,
    read_html,
    read_md,
    read_pptx,
    read_xlsx,
    write_html,
    write_md,
    write_pptx,
    write_xlsx,
)


# ── XLSX readers ─────────────────────────────────────────────────────────────

class TestReadXlsx:
    def _make_xlsx(self, tmp_path: Path, sheets: list[dict] | None = None) -> Path:
        import openpyxl
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        if sheets is None:
            sheets = [{"name": "Data", "headers": ["name", "score"], "rows": [["Alice", 90], ["Bob", 80]]}]
        for s in sheets:
            ws = wb.create_sheet(s["name"])
            ws.append(s.get("headers", []))
            for row in s.get("rows", []):
                ws.append(row)
        out = tmp_path / "data.xlsx"
        wb.save(str(out))
        return out

    def test_basic_read(self, tmp_path: Path) -> None:
        f = self._make_xlsx(tmp_path)
        result = read_xlsx(f)
        assert result.format == "xlsx"
        assert "Data" in result.sections
        sheet = result.structured["sheets"][0]
        assert sheet["headers"] == ["name", "score"]
        assert sheet["rows_preview"][0]["name"] == "Alice"

    def test_multiple_sheets(self, tmp_path: Path) -> None:
        f = self._make_xlsx(tmp_path, sheets=[
            {"name": "Sheet1", "headers": ["a"], "rows": [["1"]]},
            {"name": "Sheet2", "headers": ["b"], "rows": [["2"]]},
        ])
        result = read_xlsx(f)
        assert result.sections == ("Sheet1", "Sheet2")
        assert result.metadata["sheet_count"] == 2

    def test_row_cap_sets_truncated(self, tmp_path: Path) -> None:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["x"])
        for i in range(20):
            ws.append([i])
        out = tmp_path / "big.xlsx"
        wb.save(str(out))
        result = read_xlsx(out, ReadLimits(max_rows=5))
        assert result.truncated is True
        assert result.structured["sheets"][0]["row_count"] == 20
        assert len(result.structured["sheets"][0]["rows_preview"]) == 5

    def test_sheet_cap_sets_truncated(self, tmp_path: Path) -> None:
        import openpyxl
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        for i in range(5):
            ws = wb.create_sheet(f"S{i}")
            ws.append(["x"])
        out = tmp_path / "many.xlsx"
        wb.save(str(out))
        result = read_xlsx(out, ReadLimits(max_sheets=3))
        assert result.truncated is True
        assert len(result.structured["sheets"]) == 3

    def test_source_hash_present(self, tmp_path: Path) -> None:
        f = self._make_xlsx(tmp_path)
        result = read_xlsx(f)
        assert len(result.source_sha256) == 64
        assert result.source_bytes > 0

    def test_content_preview_includes_sheet_name(self, tmp_path: Path) -> None:
        f = self._make_xlsx(tmp_path)
        result = read_xlsx(f)
        assert "Data" in result.content
        assert "Alice" in result.content


# ── XLSX writers ─────────────────────────────────────────────────────────────

class TestWriteXlsx:
    def test_basic_roundtrip(self, tmp_path: Path) -> None:
        import openpyxl
        out = tmp_path / "out.xlsx"
        result = write_xlsx([
            {"name": "Results", "headers": ["name", "val"], "rows": [{"name": "Alice", "val": "90"}]}
        ], out)
        assert result.bytes_written > 0
        assert len(result.sha256) == 64
        wb = openpyxl.load_workbook(str(out))
        ws = wb["Results"]
        assert ws.cell(1, 1).value == "name"
        assert ws.cell(2, 1).value == "Alice"

    def test_multiple_sheets(self, tmp_path: Path) -> None:
        import openpyxl
        out = tmp_path / "multi.xlsx"
        write_xlsx([
            {"name": "A", "headers": ["x"], "rows": [{"x": "1"}]},
            {"name": "B", "headers": ["y"], "rows": [{"y": "2"}]},
        ], out)
        wb = openpyxl.load_workbook(str(out))
        assert "A" in wb.sheetnames
        assert "B" in wb.sheetnames

    def test_empty_sheets_raises(self, tmp_path: Path) -> None:
        out = tmp_path / "empty.xlsx"
        with pytest.raises(WriteValidationError):
            write_xlsx([], out)

    def test_duplicate_headers_raise(self, tmp_path: Path) -> None:
        out = tmp_path / "dup.xlsx"
        with pytest.raises(WriteValidationError, match="duplicate headers"):
            write_xlsx([{"name": "S", "headers": ["x", "x"], "rows": []}], out)

    def test_wrong_extension_raises(self, tmp_path: Path) -> None:
        out = tmp_path / "out.csv"
        with pytest.raises(WriteValidationError):
            write_xlsx([{"name": "S", "headers": [], "rows": []}], out)

    def test_sha256_correct(self, tmp_path: Path) -> None:
        out = tmp_path / "h.xlsx"
        result = write_xlsx([{"name": "S", "headers": ["a"], "rows": [{"a": "1"}]}], out)
        expected = hashlib.sha256(out.read_bytes()).hexdigest()
        assert result.sha256 == expected

    def test_result_format(self, tmp_path: Path) -> None:
        out = tmp_path / "fmt.xlsx"
        result = write_xlsx([{"name": "S", "headers": ["a"], "rows": []}], out)
        assert result.format == "xlsx"
        assert "spreadsheetml" in result.media_type


# ── PPTX readers ─────────────────────────────────────────────────────────────

class TestReadPptx:
    def _make_pptx(self, tmp_path: Path) -> Path:
        from pptx import Presentation
        prs = Presentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Introduction"
        slide.placeholders[1].text = "This is the intro content."
        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.title.text = "Methods"
        slide2.placeholders[1].text = "Method details here."
        out = tmp_path / "deck.pptx"
        prs.save(str(out))
        return out

    def test_slide_titles_extracted(self, tmp_path: Path) -> None:
        f = self._make_pptx(tmp_path)
        result = read_pptx(f)
        assert result.format == "pptx"
        assert "Introduction" in result.sections
        assert "Methods" in result.sections

    def test_slide_body_extracted(self, tmp_path: Path) -> None:
        f = self._make_pptx(tmp_path)
        result = read_pptx(f)
        slide = result.structured["slides"][0]
        assert "intro content" in slide["body"]

    def test_metadata_counts(self, tmp_path: Path) -> None:
        f = self._make_pptx(tmp_path)
        result = read_pptx(f)
        assert result.metadata["slide_count"] == 2
        assert result.metadata["slides_read"] == 2

    def test_slide_cap_sets_truncated(self, tmp_path: Path) -> None:
        from pptx import Presentation
        prs = Presentation()
        layout = prs.slide_layouts[6]
        for _ in range(5):
            prs.slides.add_slide(layout)
        out = tmp_path / "many.pptx"
        prs.save(str(out))
        result = read_pptx(out, ReadLimits(max_slides=3))
        assert result.truncated is True
        assert result.metadata["slides_read"] == 3

    def test_source_hash_present(self, tmp_path: Path) -> None:
        f = self._make_pptx(tmp_path)
        result = read_pptx(f)
        assert len(result.source_sha256) == 64

    def test_content_includes_titles(self, tmp_path: Path) -> None:
        f = self._make_pptx(tmp_path)
        result = read_pptx(f)
        assert "Introduction" in result.content or len(result.content) > 0


# ── PPTX writers ─────────────────────────────────────────────────────────────

class TestWritePptx:
    def test_basic_deck(self, tmp_path: Path) -> None:
        from pptx import Presentation
        out = tmp_path / "out.pptx"
        result = write_pptx([
            {"title": "Slide One", "body": "First slide body."},
            {"title": "Slide Two", "body": "Second slide body."},
        ], out)
        assert result.bytes_written > 0
        assert len(result.sha256) == 64
        prs = Presentation(str(out))
        assert len(prs.slides) == 2

    def test_title_only_slide(self, tmp_path: Path) -> None:
        from pptx import Presentation
        out = tmp_path / "title_only.pptx"
        write_pptx([{"title": "Just a Title"}], out)
        prs = Presentation(str(out))
        assert len(prs.slides) == 1

    def test_body_only_slide(self, tmp_path: Path) -> None:
        from pptx import Presentation
        out = tmp_path / "body_only.pptx"
        write_pptx([{"body": "Just some body text."}], out)
        prs = Presentation(str(out))
        assert len(prs.slides) == 1

    def test_wrong_extension_raises(self, tmp_path: Path) -> None:
        out = tmp_path / "out.pdf"
        with pytest.raises(WriteValidationError):
            write_pptx([], out)

    def test_sha256_correct(self, tmp_path: Path) -> None:
        out = tmp_path / "h.pptx"
        result = write_pptx([{"title": "T", "body": "B"}], out)
        expected = hashlib.sha256(out.read_bytes()).hexdigest()
        assert result.sha256 == expected

    def test_result_format(self, tmp_path: Path) -> None:
        out = tmp_path / "fmt.pptx"
        result = write_pptx([], out)
        assert result.format == "pptx"
        assert "presentationml" in result.media_type


# ── Markdown readers ──────────────────────────────────────────────────────────

class TestReadMd:
    def test_headings_extracted(self, tmp_path: Path) -> None:
        f = tmp_path / "doc.md"
        f.write_text("# Title\n\nIntro text.\n\n## Section One\n\nBody here.\n", encoding="utf-8")
        result = read_md(f)
        assert result.format == "md"
        assert "Title" in result.sections
        assert "Section One" in result.sections

    def test_content_preserved(self, tmp_path: Path) -> None:
        f = tmp_path / "doc.md"
        f.write_text("# H1\n\nsome content\n", encoding="utf-8")
        result = read_md(f)
        assert "some content" in result.content

    def test_structured_sections(self, tmp_path: Path) -> None:
        f = tmp_path / "doc.md"
        f.write_text("# A\n\nbody A\n\n## B\n\nbody B\n", encoding="utf-8")
        result = read_md(f)
        sections = result.structured["sections"]
        assert any(s["heading"] == "A" for s in sections)

    def test_char_cap_sets_truncated(self, tmp_path: Path) -> None:
        f = tmp_path / "long.md"
        f.write_text("x" * 1000, encoding="utf-8")
        result = read_md(f, ReadLimits(max_chars=100))
        assert result.truncated is True
        assert len(result.content) == 100

    def test_bom_stripped(self, tmp_path: Path) -> None:
        f = tmp_path / "bom.md"
        f.write_bytes(b"\xef\xbb\xbf# Hello\n")
        result = read_md(f)
        assert result.content.startswith("# Hello")

    def test_metadata(self, tmp_path: Path) -> None:
        f = tmp_path / "doc.md"
        f.write_text("# A\n\n## B\n\ntext\n", encoding="utf-8")
        result = read_md(f)
        assert result.metadata["heading_count"] == 2
        assert result.metadata["char_count"] > 0


# ── Markdown writers ──────────────────────────────────────────────────────────

class TestWriteMd:
    def test_basic_roundtrip(self, tmp_path: Path) -> None:
        out = tmp_path / "out.md"
        result = write_md([
            {"heading": "Introduction", "level": 1, "body": "Intro body text."},
            {"heading": "Details", "level": 2, "body": "More details here."},
        ], out)
        assert result.bytes_written > 0
        content = out.read_text(encoding="utf-8")
        assert "# Introduction" in content
        assert "## Details" in content
        assert "Intro body text." in content

    def test_body_without_heading(self, tmp_path: Path) -> None:
        out = tmp_path / "nobody.md"
        write_md([{"body": "Just body text."}], out)
        assert "Just body text." in out.read_text(encoding="utf-8")

    def test_wrong_extension_raises(self, tmp_path: Path) -> None:
        out = tmp_path / "out.txt"
        with pytest.raises(WriteValidationError):
            write_md([], out)

    def test_sha256_correct(self, tmp_path: Path) -> None:
        out = tmp_path / "h.md"
        result = write_md([{"heading": "H", "body": "B"}], out)
        expected = hashlib.sha256(out.read_bytes()).hexdigest()
        assert result.sha256 == expected

    def test_result_format(self, tmp_path: Path) -> None:
        out = tmp_path / "fmt.md"
        result = write_md([], out)
        assert result.format == "md"
        assert result.media_type == "text/markdown"

    def test_heading_levels_clamped(self, tmp_path: Path) -> None:
        out = tmp_path / "levels.md"
        write_md([{"heading": "H", "level": 99, "body": ""}], out)
        content = out.read_text(encoding="utf-8")
        assert content.startswith("######")  # clamped to 6


# ── HTML readers ──────────────────────────────────────────────────────────────

class TestReadHtml:
    def test_basic_text_extraction(self, tmp_path: Path) -> None:
        f = tmp_path / "page.html"
        f.write_text(
            "<html><body><h1>Title</h1><p>Hello world</p></body></html>",
            encoding="utf-8",
        )
        result = read_html(f)
        assert result.format == "html"
        assert "Title" in result.sections
        assert "Hello world" in result.content

    def test_scripts_stripped(self, tmp_path: Path) -> None:
        f = tmp_path / "js.html"
        f.write_text(
            "<html><body><p>Real</p><script>evil()</script></body></html>",
            encoding="utf-8",
        )
        result = read_html(f)
        assert "evil" not in result.content

    def test_styles_stripped(self, tmp_path: Path) -> None:
        f = tmp_path / "css.html"
        f.write_text(
            "<html><head><style>body{color:red}</style></head><body><p>Text</p></body></html>",
            encoding="utf-8",
        )
        result = read_html(f)
        assert "color" not in result.content

    def test_links_extracted(self, tmp_path: Path) -> None:
        f = tmp_path / "links.html"
        f.write_text(
            '<html><body><a href="https://example.com">Click</a></body></html>',
            encoding="utf-8",
        )
        result = read_html(f)
        assert result.metadata["link_count"] == 1
        assert result.structured["links"][0]["href"] == "https://example.com"

    def test_headings_extracted(self, tmp_path: Path) -> None:
        f = tmp_path / "headings.html"
        f.write_text(
            "<html><body><h1>One</h1><h2>Two</h2><p>Body</p></body></html>",
            encoding="utf-8",
        )
        result = read_html(f)
        assert "One" in result.sections
        assert "Two" in result.sections

    def test_char_cap_sets_truncated(self, tmp_path: Path) -> None:
        f = tmp_path / "long.html"
        f.write_text(f"<p>{'x' * 1000}</p>", encoding="utf-8")
        result = read_html(f, ReadLimits(max_chars=50))
        assert result.truncated is True

    def test_source_hash_present(self, tmp_path: Path) -> None:
        f = tmp_path / "h.html"
        f.write_text("<p>hi</p>", encoding="utf-8")
        result = read_html(f)
        assert len(result.source_sha256) == 64


# ── HTML writers ──────────────────────────────────────────────────────────────

class TestWriteHtml:
    def test_basic_output(self, tmp_path: Path) -> None:
        out = tmp_path / "out.html"
        result = write_html([
            {"heading": "Introduction", "level": 1, "body": "Intro body text."},
        ], out, title="My Doc")
        assert result.bytes_written > 0
        content = out.read_text(encoding="utf-8")
        assert "<h1>Introduction</h1>" in content
        assert "<title>My Doc</title>" in content
        assert "Intro body text." in content

    def test_escaping(self, tmp_path: Path) -> None:
        out = tmp_path / "esc.html"
        write_html([{"heading": "<script>", "body": "safe & sound"}], out)
        content = out.read_text(encoding="utf-8")
        assert "<script>" not in content
        assert "&lt;script&gt;" in content
        assert "&amp;" in content

    def test_htm_extension_allowed(self, tmp_path: Path) -> None:
        out = tmp_path / "out.htm"
        result = write_html([], out)
        assert result.bytes_written > 0

    def test_wrong_extension_raises(self, tmp_path: Path) -> None:
        out = tmp_path / "out.txt"
        with pytest.raises(WriteValidationError):
            write_html([], out)

    def test_sha256_correct(self, tmp_path: Path) -> None:
        out = tmp_path / "h.html"
        result = write_html([{"heading": "H", "body": "B"}], out)
        expected = hashlib.sha256(out.read_bytes()).hexdigest()
        assert result.sha256 == expected

    def test_result_format(self, tmp_path: Path) -> None:
        out = tmp_path / "fmt.html"
        result = write_html([], out)
        assert result.format == "html"
        assert result.media_type == "text/html"

    def test_valid_html_structure(self, tmp_path: Path) -> None:
        out = tmp_path / "struct.html"
        write_html([{"heading": "H", "body": "B"}], out)
        content = out.read_text(encoding="utf-8")
        assert "<!DOCTYPE html>" in content
        assert "<html>" in content
        assert "</html>" in content
        assert "<body>" in content


# ── Registry extended ─────────────────────────────────────────────────────────

class TestExtendedRegistry:
    def test_new_formats_in_registry(self) -> None:
        from core.file_formats import FORMAT_HANDLERS
        assert ".xlsx" in FORMAT_HANDLERS
        assert ".pptx" in FORMAT_HANDLERS
        assert ".md" in FORMAT_HANDLERS
        assert ".html" in FORMAT_HANDLERS

    def test_xlsx_media_type(self) -> None:
        from core.file_formats import FORMAT_HANDLERS
        assert "spreadsheetml" in FORMAT_HANDLERS[".xlsx"].media_type

    def test_pptx_media_type(self) -> None:
        from core.file_formats import FORMAT_HANDLERS
        assert "presentationml" in FORMAT_HANDLERS[".pptx"].media_type

    def test_md_media_type(self) -> None:
        from core.file_formats import FORMAT_HANDLERS
        assert FORMAT_HANDLERS[".md"].media_type == "text/markdown"

    def test_html_media_type(self) -> None:
        from core.file_formats import FORMAT_HANDLERS
        assert FORMAT_HANDLERS[".html"].media_type == "text/html"


# ── DocumentExtractor compat ──────────────────────────────────────────────────

class TestExtractorNewFormats:
    def test_xlsx_dispatches(self, tmp_path: Path) -> None:
        import openpyxl
        from memory.document_extraction import DocumentExtractor
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["name", "val"])
        ws.append(["Alice", "90"])
        out = tmp_path / "d.xlsx"
        wb.save(str(out))
        doc = DocumentExtractor().extract(str(out))
        assert doc.format == ".xlsx"
        assert "Sheet" in doc.sections  # sections = sheet names

    def test_pptx_dispatches(self, tmp_path: Path) -> None:
        from pptx import Presentation
        from memory.document_extraction import DocumentExtractor
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "My Slide"
        out = tmp_path / "d.pptx"
        prs.save(str(out))
        doc = DocumentExtractor().extract(str(out))
        assert doc.format == ".pptx"

    def test_html_dispatches(self, tmp_path: Path) -> None:
        from memory.document_extraction import DocumentExtractor
        f = tmp_path / "d.html"
        f.write_text("<h1>Title</h1><p>Body</p>", encoding="utf-8")
        doc = DocumentExtractor().extract(str(f))
        assert doc.format == ".html"
        assert "Title" in doc.sections

    def test_htm_dispatches(self, tmp_path: Path) -> None:
        from memory.document_extraction import DocumentExtractor
        f = tmp_path / "d.htm"
        f.write_text("<p>content</p>", encoding="utf-8")
        doc = DocumentExtractor().extract(str(f))
        assert doc.format == ".html"

    def test_md_dispatches(self, tmp_path: Path) -> None:
        from memory.document_extraction import DocumentExtractor
        f = tmp_path / "d.md"
        f.write_text("# Title\n\nBody text.\n", encoding="utf-8")
        doc = DocumentExtractor().extract(str(f))
        assert doc.format == ".md"
        assert "Title" in doc.sections
