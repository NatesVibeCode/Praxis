"""Deterministic validators for Model Eval task outputs."""

from __future__ import annotations

import csv
import base64
import io
import json
from typing import Any


def _artifact(payload: dict[str, Any], suffix: str) -> dict[str, str] | None:
    for item in payload.get("artifacts") or []:
        if not isinstance(item, dict):
            continue
        path = str(item.get("path") or "")
        if path.endswith(suffix):
            return {
                "path": path,
                "media_type": str(item.get("media_type") or ""),
                "content": str(item.get("content") or ""),
            }
    return None


def _check(ok: bool, check: str, detail: Any = None) -> dict[str, Any]:
    return {"ok": bool(ok), "check": check, "detail": detail}


def _csv_rows(text: str) -> list[dict[str, Any]]:
    return list(csv.DictReader(io.StringIO(text)))


def _stringify_cell(value: Any) -> str:
    if isinstance(value, list):
        return " ".join(str(item) for item in value)
    if value is None:
        return ""
    return str(value)


def _row_text(row: dict[Any, Any]) -> str:
    return " ".join(_stringify_cell(value) for value in row.values())


def validate_doc_user_guide(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    artifact = _artifact(payload, "user_guide.md")
    text = artifact["content"] if artifact else ""
    headings = [
        "# Model Eval User Guide",
        "## What It Tests",
        "## Running A Matrix",
        "## Reading Results",
        "## Promotion Rules",
        "## Troubleshooting",
    ]
    checks = [_check(artifact is not None, "user_guide.md artifact exists")]
    checks.extend(_check(heading in text, f"heading {heading}") for heading in headings)
    checks.append(_check("production routing" in text.lower(), "production routing boundary"))
    checks.append(_check("promot" in text.lower(), "promotion rule mentioned"))
    checks.append(_check(600 <= len(text) <= 5000, "guide length bounded", len(text)))
    return _result(checks)


def validate_structured_doc_headings(payload: dict[str, Any], task: dict[str, Any]) -> dict[str, Any]:
    expected_path = str(task.get("expected_artifact") or ".md")
    artifact = _artifact(payload, expected_path)
    text = artifact["content"] if artifact else ""
    headings = [str(item) for item in (task.get("expected_headings") or []) if str(item).strip()]
    checks = [_check(artifact is not None, f"{expected_path} artifact exists")]
    checks.extend(_check(heading in text, f"heading {heading}") for heading in headings)
    checks.append(_check("unknown" in text.lower() or "assumption" in text.lower(), "unknowns or assumptions named"))
    checks.append(_check("production routing" in text.lower() or "routing unchanged" in text.lower(), "routing boundary"))
    checks.append(_check(500 <= len(text) <= 6000, "document length bounded", len(text)))
    return _result(checks)


def validate_pptx_deck_manifest(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    artifact = _artifact(payload, "deck.json")
    checks = [_check(artifact is not None, "deck.json artifact exists")]
    document: dict[str, Any] = {}
    if artifact:
        try:
            parsed = json.loads(artifact["content"])
            document = parsed if isinstance(parsed, dict) else {}
            checks.append(_check(isinstance(parsed, dict), "deck.json parses as object"))
        except json.JSONDecodeError as exc:
            checks.append(_check(False, "deck.json parses", str(exc)))
    slides = document.get("slides")
    checks.append(_check(isinstance(slides, list) and len(slides) == 6, "exactly six slides"))
    if isinstance(slides, list):
        required_topics = ("privacy", "task", "cost", "consistency", "promotion")
        blob = json.dumps(slides).lower()
        for topic in required_topics:
            checks.append(_check(topic in blob, f"topic {topic} covered"))
        for index, slide in enumerate(slides):
            if not isinstance(slide, dict):
                checks.append(_check(False, f"slide {index + 1} object"))
                continue
            checks.append(_check(bool(slide.get("title")), f"slide {index + 1} title"))
            checks.append(_check(isinstance(slide.get("bullets"), list), f"slide {index + 1} bullets"))
            checks.append(_check(bool(slide.get("speaker_notes")), f"slide {index + 1} notes"))
    return _result(checks)


def validate_pptx_render(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    artifact = _artifact(payload, ".pptx")
    checks = [_check(artifact is not None, ".pptx artifact exists")]
    if not artifact:
        return _result(checks)
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # noqa: BLE001 - optional runtime dependency.
        checks.append(_check(False, "python-pptx available", f"{type(exc).__name__}: {exc}"))
        return _result(checks)
    try:
        raw = base64.b64decode(artifact["content"], validate=True)
        deck = Presentation(io.BytesIO(raw))
        checks.append(_check(True, ".pptx opens"))
    except Exception as exc:  # noqa: BLE001 - malformed deck is verifier data.
        checks.append(_check(False, ".pptx opens", f"{type(exc).__name__}: {exc}"))
        return _result(checks)
    checks.append(_check(len(deck.slides) >= 4, "at least four slides", len(deck.slides)))
    titles = []
    for slide in deck.slides:
        title = getattr(slide.shapes, "title", None)
        if title is not None and getattr(title, "text", ""):
            titles.append(title.text.strip())
    checks.append(_check(len(titles) >= min(4, len(deck.slides)), "slide titles present", titles))
    return _result(checks)


def validate_csv_extract_accounts(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    artifact = _artifact(payload, "extracted_accounts.csv")
    checks = [_check(artifact is not None, "extracted_accounts.csv artifact exists")]
    rows: list[dict[str, Any]] = []
    if artifact:
        rows = _csv_rows(artifact["content"])
    expected_columns = ["account_id", "owner", "status", "next_action", "risk_score"]
    checks.append(_check(rows and list(rows[0].keys()) == expected_columns, "exact columns"))
    checks.append(_check(all(None not in row for row in rows), "no CSV overflow columns"))
    checks.append(_check(len(rows) == 4, "four rows", len(rows)))
    by_id = {row.get("account_id"): row for row in rows}
    checks.append(_check(by_id.get("A-17", {}).get("risk_score") == "91", "A-17 risk score"))
    checks.append(_check("credential, then retry" in by_id.get("A-17", {}).get("next_action", ""), "quoted comma preserved"))
    return _result(checks)


def validate_csv_create_rollout(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    artifact = _artifact(payload, "rollout_plan.csv")
    checks = [_check(artifact is not None, "rollout_plan.csv artifact exists")]
    rows: list[dict[str, Any]] = []
    if artifact:
        rows = _csv_rows(artifact["content"])
    expected_columns = ["week", "workstream", "owner", "deliverable", "done_definition"]
    checks.append(_check(rows and list(rows[0].keys()) == expected_columns, "exact columns"))
    checks.append(_check(all(None not in row for row in rows), "no CSV overflow columns"))
    checks.append(_check(len(rows) == 5, "five rows", len(rows)))
    blob = " ".join(_row_text(row).lower() for row in rows)
    for term in ("docs", "pptx", "csv", "tool", "swarm"):
        checks.append(_check(term in blob, f"covers {term}"))
    return _result(checks)


def validate_csv_reconcile_accounts(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    artifact = _artifact(payload, "account_reconciliation.csv")
    checks = [_check(artifact is not None, "account_reconciliation.csv artifact exists")]
    rows: list[dict[str, Any]] = []
    if artifact:
        rows = _csv_rows(artifact["content"])
    expected_columns = ["account_id", "source_a_status", "source_b_status", "disposition", "notes"]
    checks.append(_check(rows and list(rows[0].keys()) == expected_columns, "exact columns"))
    checks.append(_check(all(None not in row for row in rows), "no CSV overflow columns"))
    checks.append(_check(len(rows) == 4, "four rows", len(rows)))
    by_id = {row.get("account_id"): row for row in rows}
    checks.append(_check(by_id.get("A-17", {}).get("disposition") == "conflict", "A-17 conflict"))
    checks.append(_check(by_id.get("B-04", {}).get("disposition") == "match", "B-04 match"))
    checks.append(_check("review" in by_id.get("C-88", {}).get("notes", "").lower(), "C-88 review note"))
    return _result(checks)


def validate_workbook_manifest(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    artifact = _artifact(payload, "workbook_manifest.json")
    checks = [_check(artifact is not None, "workbook_manifest.json artifact exists")]
    document: dict[str, Any] = {}
    if artifact:
        try:
            parsed = json.loads(artifact["content"])
            document = parsed if isinstance(parsed, dict) else {}
            checks.append(_check(isinstance(parsed, dict), "workbook manifest parses as object"))
        except json.JSONDecodeError as exc:
            checks.append(_check(False, "workbook manifest parses", str(exc)))
    sheets = document.get("sheets")
    checks.append(_check(isinstance(sheets, list) and len(sheets) >= 2, "at least two sheets"))
    blob = json.dumps(document).lower()
    checks.append(_check("formula" in blob or "=" in blob, "formulas represented"))
    checks.append(_check("chart" in blob, "chart represented"))
    checks.append(_check("recalc" in blob or "calculation" in blob, "calculation policy represented"))
    return _result(checks)


def _tool_call_names(payload: dict[str, Any]) -> list[str]:
    calls = payload.get("tool_calls") or []
    names: list[str] = []
    for call in calls:
        if not isinstance(call, dict):
            continue
        function = call.get("function") if isinstance(call.get("function"), dict) else {}
        names.append(str(function.get("name") or call.get("name") or ""))
    return names


def _validate_tool_names(payload: dict[str, Any], expected: list[str]) -> dict[str, Any]:
    names = _tool_call_names(payload)
    checks = [
        _check(len(names) == len(expected), f"exactly {len(expected)} tool call(s)", names),
        _check(names == expected, "expected tool order", names),
    ]
    for name in expected:
        checks.append(_check(name in names, f"calls {name}", names))
    return _result(checks)


def validate_tool_single_search(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    return _validate_tool_names(payload, ["praxis_search"])


def validate_tool_single_validate(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    return _validate_tool_names(payload, ["praxis_workflow_validate"])


def validate_tool_single_model_eval(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    return _validate_tool_names(payload, ["praxis_model_eval"])


def validate_tool_single_bugs(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    return _validate_tool_names(payload, ["praxis_bugs"])


def validate_tool_single_operator_decisions(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    return _validate_tool_names(payload, ["praxis_operator_decisions"])


def validate_tool_choice_search(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    return _validate_tool_names(payload, ["praxis_search"])


def validate_tool_choice_model_eval(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    return _validate_tool_names(payload, ["praxis_model_eval"])


def validate_tool_call_sequence(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    names = _tool_call_names(payload)
    checks = [
        _check(len(names) == 2, "exactly two tool calls", names),
        _check(names[:1] == ["praxis_search"], "first tool is search", names),
        _check(names[1:2] == ["praxis_workflow_validate"], "second tool is validation", names),
    ]
    return _result(checks)


def validate_tool_execution_transcript(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    artifact = _artifact(payload, "tool_transcript.json")
    checks = [_check(artifact is not None, "tool_transcript.json artifact exists")]
    transcript: dict[str, Any] = {}
    if artifact:
        try:
            parsed = json.loads(artifact["content"])
            transcript = parsed if isinstance(parsed, dict) else {}
            checks.append(_check(isinstance(parsed, dict), "tool transcript parses as object"))
        except json.JSONDecodeError as exc:
            checks.append(_check(False, "tool transcript parses", str(exc)))
    steps = transcript.get("steps")
    checks.append(_check(isinstance(steps, list) and len(steps) >= 2, "multi-step transcript"))
    if isinstance(steps, list):
        tool_calls = [
            step
            for step in steps
            if isinstance(step, dict) and step.get("kind") == "tool_call"
        ]
        tool_results = [
            step
            for step in steps
            if isinstance(step, dict) and step.get("kind") == "tool_result"
        ]
        checks.append(_check(bool(tool_calls), "tool call recorded"))
        checks.append(_check(bool(tool_results), "tool result recorded"))
        names = [str(step.get("tool_name") or "") for step in tool_calls]
        checks.append(_check("praxis_search" in names, "search tool used", names))
        receipt_refs = [step.get("receipt_id") for step in tool_results if step.get("receipt_id")]
        checks.append(_check(bool(receipt_refs) or any(step.get("ok") is False for step in tool_results), "receipt or explicit refusal recorded"))
    return _result(checks)


def validate_swarm_packet(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    artifact = _artifact(payload, "swarm_plan.json")
    checks = [_check(artifact is not None, "swarm_plan.json artifact exists")]
    document: dict[str, Any] = {}
    if artifact:
        try:
            parsed = json.loads(artifact["content"])
            document = parsed if isinstance(parsed, dict) else {}
            checks.append(_check(isinstance(parsed, dict), "swarm_plan.json parses as object"))
        except json.JSONDecodeError as exc:
            checks.append(_check(False, "swarm_plan.json parses", str(exc)))
    workers = document.get("workers")
    reducer = document.get("reducer")
    checks.append(_check(isinstance(workers, list) and len(workers) == 4, "four workers"))
    checks.append(_check(isinstance(reducer, dict) or isinstance(reducer, str), "reducer defined"))
    blob = json.dumps(document).lower()
    checks.append(_check("production routing" in blob, "production routing boundary"))
    checks.append(_check("budget" in blob, "budget cap mentioned"))
    if isinstance(workers, list):
        names = [str(item.get("name") or item.get("role") or "") for item in workers if isinstance(item, dict)]
        checks.append(_check(len(set(names)) == 4, "worker names distinct", names))
    return _result(checks)


def validate_swarm_reducer_packet(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    artifact = _artifact(payload, "swarm_reducer.json")
    checks = [_check(artifact is not None, "swarm_reducer.json artifact exists")]
    document: dict[str, Any] = {}
    if artifact:
        try:
            parsed = json.loads(artifact["content"])
            document = parsed if isinstance(parsed, dict) else {}
            checks.append(_check(isinstance(parsed, dict), "swarm_reducer.json parses as object"))
        except json.JSONDecodeError as exc:
            checks.append(_check(False, "swarm_reducer.json parses", str(exc)))
    blob = json.dumps(document).lower()
    checks.append(_check("worker_outputs" in document or "workers" in document, "worker outputs represented"))
    checks.append(_check("overlap" in blob, "overlap detection represented"))
    checks.append(_check("budget" in blob, "budget represented"))
    checks.append(_check("production routing" in blob or "routing unchanged" in blob, "routing boundary"))
    checks.append(_check("decision" in blob or "winner" in blob, "reducer decision represented"))
    return _result(checks)


def validate_workflow_job_packet(payload: dict[str, Any], _task: dict[str, Any]) -> dict[str, Any]:
    artifact = _artifact(payload, "decision_packet.md")
    text = artifact["content"] if artifact else ""
    headings = [
        "# Workflow Job Packet",
        "## Imported Spec",
        "## Proposed Work",
        "## Acceptance Evidence",
        "## Risks",
        "## Verifier Notes",
    ]
    checks = [_check(artifact is not None, "decision_packet.md artifact exists")]
    checks.extend(_check(heading in text, f"heading {heading}") for heading in headings)
    checks.append(_check("executed" not in text.lower(), "does not claim execution"))
    return _result(checks)


def _result(checks: list[dict[str, Any]]) -> dict[str, Any]:
    passed = sum(1 for item in checks if item.get("ok"))
    score = passed / max(1, len(checks))
    return {"ok": all(item.get("ok") for item in checks), "score": round(score, 4), "checks": checks}


VALIDATORS = {
    "doc_user_guide": validate_doc_user_guide,
    "structured_doc_headings": validate_structured_doc_headings,
    "pptx_deck_manifest": validate_pptx_deck_manifest,
    "pptx_render": validate_pptx_render,
    "csv_extract_accounts": validate_csv_extract_accounts,
    "csv_create_rollout": validate_csv_create_rollout,
    "csv_reconcile_accounts": validate_csv_reconcile_accounts,
    "workbook_manifest": validate_workbook_manifest,
    "tool_single_search": validate_tool_single_search,
    "tool_single_validate": validate_tool_single_validate,
    "tool_single_model_eval": validate_tool_single_model_eval,
    "tool_single_bugs": validate_tool_single_bugs,
    "tool_single_operator_decisions": validate_tool_single_operator_decisions,
    "tool_choice_search": validate_tool_choice_search,
    "tool_choice_model_eval": validate_tool_choice_model_eval,
    "tool_call_sequence": validate_tool_call_sequence,
    "tool_execution_transcript": validate_tool_execution_transcript,
    "swarm_packet": validate_swarm_packet,
    "swarm_reducer_packet": validate_swarm_reducer_packet,
    "workflow_job_packet": validate_workflow_job_packet,
}


def validate_task_output(task: dict[str, Any], payload: dict[str, Any]) -> dict[str, Any]:
    validator_name = str(task.get("validator") or "")
    validator = VALIDATORS.get(validator_name)
    if validator is None:
        return {
            "ok": False,
            "score": 0.0,
            "checks": [_check(False, "validator registered", validator_name)],
        }
    return validator(payload, task)


__all__ = ["VALIDATORS", "validate_task_output"]
