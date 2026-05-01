"""File format readers: CSV, TXT, PDF, DOCX, XLSX, PPTX, Markdown, HTML."""

from __future__ import annotations

import csv
import hashlib
import io
import os
from pathlib import Path
from typing import TYPE_CHECKING

from core.file_formats.models import (
    EncryptedDocumentError,
    FileFormatError,
    FileParseError,
    FileReadError,
    FORMAT_HANDLERS,
    ReadLimits,
    ReadResult,
    UnsupportedFileFormatError,
)

if TYPE_CHECKING:
    pass


def read_file(path: str | Path, limits: ReadLimits | None = None) -> ReadResult:
    """Dispatch to the appropriate reader based on file extension."""
    ext = os.path.splitext(str(path))[1].lower()
    handler = FORMAT_HANDLERS.get(ext)
    if handler is None:
        raise UnsupportedFileFormatError(f"No reader for extension: {ext!r}")
    _readers = {
        ".csv": read_csv,
        ".txt": read_txt,
        ".pdf": read_pdf,
        ".docx": read_docx,
        ".xlsx": read_xlsx,
        ".pptx": read_pptx,
        ".md": read_md,
        ".html": read_html,
        ".htm": read_html,
    }
    return _readers[ext](path, limits)


def read_csv(path: str | Path, limits: ReadLimits | None = None) -> ReadResult:
    lim = limits or ReadLimits()
    path = str(path)
    try:
        raw = _read_bytes(path, lim.max_bytes)
    except OSError as exc:
        raise FileReadError(str(exc)) from exc

    sha256 = _sha256(raw)
    source_bytes = len(raw)

    try:
        text = raw.decode("utf-8-sig")
    except UnicodeDecodeError:
        text = raw.decode("latin-1")

    if not text.strip():
        return ReadResult(
            path=path,
            format="csv",
            media_type="text/csv",
            content="",
            sections=(),
            metadata={"column_count": 0, "preview_row_count": 0},
            structured={"headers": [], "rows_preview": [], "delimiter": ","},
            warnings=("empty file",),
            source_sha256=sha256,
            source_bytes=source_bytes,
        )

    # Sniff dialect
    sample = text[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample)
    except csv.Error:
        dialect = csv.excel

    warnings: list[str] = []
    reader = csv.DictReader(io.StringIO(text, newline=""), dialect=dialect)

    headers = list(reader.fieldnames or [])
    if not headers:
        warnings.append("no headers detected")

    # Duplicate header check
    seen: set[str] = set()
    for h in headers:
        if h in seen:
            raise FileParseError(f"Duplicate header: {h!r}")
        seen.add(h)

    rows_preview: list[dict] = []
    truncated = False
    for row in reader:
        # Skip fully blank rows
        if all(not (v or "").strip() for v in row.values()):
            continue
        if len(rows_preview) >= lim.max_rows:
            truncated = True
            break
        # Warn on field count mismatches (DictReader uses None key for extras)
        if None in row:
            warnings.append(f"row {len(rows_preview) + 1} has more fields than headers")
            row = {k: v for k, v in row.items() if k is not None}
        cleaned = {h: row.get(h) or "" for h in headers}
        rows_preview.append(cleaned)

    # Build content preview
    content_lines = [",".join(headers)]
    for row in rows_preview:
        content_lines.append(",".join(str(row.get(h, "")) for h in headers))
    content = "\n".join(content_lines)
    if truncated:
        content += f"\n... (truncated at {lim.max_rows} rows)"

    return ReadResult(
        path=path,
        format="csv",
        media_type="text/csv",
        content=content,
        sections=tuple(headers),
        metadata={
            "column_count": len(headers),
            "preview_row_count": len(rows_preview),
        },
        structured={
            "headers": headers,
            "rows_preview": rows_preview,
            "delimiter": getattr(dialect, "delimiter", ","),
        },
        warnings=tuple(warnings),
        truncated=truncated,
        source_sha256=sha256,
        source_bytes=source_bytes,
    )


def read_txt(path: str | Path, limits: ReadLimits | None = None) -> ReadResult:
    lim = limits or ReadLimits()
    path = str(path)
    try:
        raw = _read_bytes(path, lim.max_bytes)
    except OSError as exc:
        raise FileReadError(str(exc)) from exc

    sha256 = _sha256(raw)
    source_bytes = len(raw)

    warnings: list[str] = []
    try:
        text = raw.decode("utf-8-sig")
    except UnicodeDecodeError:
        text = raw.decode("latin-1", errors="replace")
        warnings.append("non-UTF-8 characters replaced")

    text = text.replace("\r\n", "\n").replace("\r", "\n")

    truncated = False
    if len(text) > lim.max_chars:
        text = text[: lim.max_chars]
        truncated = True
        warnings.append(f"content truncated at {lim.max_chars} characters")

    return ReadResult(
        path=path,
        format="txt",
        media_type="text/plain",
        content=text,
        sections=("Text",),
        metadata={
            "line_count": text.count("\n") + 1 if text else 0,
            "char_count": len(text),
        },
        structured={},
        warnings=tuple(warnings),
        truncated=truncated,
        source_sha256=sha256,
        source_bytes=source_bytes,
    )


def read_pdf(path: str | Path, limits: ReadLimits | None = None) -> ReadResult:
    lim = limits or ReadLimits()
    path = str(path)
    try:
        import pypdf
    except ImportError as exc:
        raise FileFormatError("pypdf is required for PDF reading") from exc

    try:
        raw = _read_bytes(path, lim.max_bytes)
    except OSError as exc:
        raise FileReadError(str(exc)) from exc

    sha256 = _sha256(raw)
    source_bytes = len(raw)

    try:
        reader = pypdf.PdfReader(io.BytesIO(raw))
    except Exception as exc:
        raise FileReadError(f"Could not open PDF: {exc}") from exc

    if reader.is_encrypted:
        raise EncryptedDocumentError(
            "PDF is encrypted; password-protected PDFs are not supported in v1"
        )

    total_pages = len(reader.pages)
    warnings: list[str] = []
    pages_data: list[dict] = []
    content_parts: list[str] = []
    total_chars = 0
    truncated = False

    for i, page in enumerate(reader.pages):
        if i >= lim.max_pages:
            truncated = True
            break
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
            warnings.append(f"page {i + 1} text extraction failed")

        if not text.strip():
            warnings.append(f"page {i + 1} returned no text (may be image-only)")

        total_chars += len(text)
        if total_chars > lim.max_chars:
            # Trim last page to fit
            overflow = total_chars - lim.max_chars
            text = text[: len(text) - overflow]
            truncated = True

        pages_data.append({"page": i + 1, "text": text, "char_count": len(text)})
        content_parts.append(f"--- Page {i + 1} ---\n{text}")

        if truncated:
            break

    content = "\n".join(content_parts)
    sections = tuple(f"Page {p['page']}" for p in pages_data)

    return ReadResult(
        path=path,
        format="pdf",
        media_type="application/pdf",
        content=content,
        sections=sections,
        metadata={
            "page_count": total_pages,
            "pages_read": len(pages_data),
            "encrypted": False,
        },
        structured={"pages": pages_data},
        warnings=tuple(warnings),
        truncated=truncated,
        source_sha256=sha256,
        source_bytes=source_bytes,
    )


def read_docx(path: str | Path, limits: ReadLimits | None = None) -> ReadResult:
    lim = limits or ReadLimits()
    path = str(path)
    try:
        import docx as python_docx
    except ImportError as exc:
        raise FileFormatError("python-docx is required for DOCX reading") from exc

    try:
        raw = _read_bytes(path, lim.max_bytes)
    except OSError as exc:
        raise FileReadError(str(exc)) from exc

    sha256 = _sha256(raw)
    source_bytes = len(raw)

    try:
        doc = python_docx.Document(io.BytesIO(raw))
    except Exception as exc:
        raise FileReadError(f"Could not open DOCX: {exc}") from exc

    warnings: list[str] = []
    headings: list[dict] = []
    paragraphs_data: list[dict] = []
    content_parts: list[str] = []
    truncated = False
    total_chars = 0

    for idx, para in enumerate(doc.paragraphs):
        if len(paragraphs_data) >= lim.max_paragraphs:
            truncated = True
            break
        text = para.text
        style_name = para.style.name if para.style else "Normal"

        if style_name.startswith("Heading"):
            headings.append({
                "text": text,
                "style": style_name,
                "paragraph_index": idx,
            })

        paragraphs_data.append({"index": idx, "text": text, "style": style_name})
        content_parts.append(text)
        total_chars += len(text)

        if total_chars > lim.max_chars:
            truncated = True
            break

    # Extract tables
    tables_data: list[dict] = []
    for t_idx, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            rows.append([cell.text for cell in row.cells])
        tables_data.append({"index": t_idx, "rows": rows})
        # Append table preview to content
        for row in rows:
            content_parts.append(" | ".join(row))

    content = "\n".join(content_parts)
    sections = tuple(h["text"] for h in headings)

    return ReadResult(
        path=path,
        format="docx",
        media_type=(
            "application/vnd.openxmlformats-officedocument"
            ".wordprocessingml.document"
        ),
        content=content,
        sections=sections,
        metadata={
            "paragraph_count": len(paragraphs_data),
            "heading_count": len(headings),
            "table_count": len(tables_data),
        },
        structured={
            "headings": headings,
            "paragraphs": paragraphs_data,
            "tables": tables_data,
        },
        warnings=tuple(warnings),
        truncated=truncated,
        source_sha256=sha256,
        source_bytes=source_bytes,
    )


def read_xlsx(path: str | Path, limits: ReadLimits | None = None) -> ReadResult:
    lim = limits or ReadLimits()
    path = str(path)
    try:
        import openpyxl
    except ImportError as exc:
        raise FileFormatError("openpyxl is required for XLSX reading") from exc

    try:
        raw = _read_bytes(path, lim.max_bytes)
    except OSError as exc:
        raise FileReadError(str(exc)) from exc

    sha256 = _sha256(raw)
    source_bytes = len(raw)

    try:
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    except Exception as exc:
        raise FileReadError(f"Could not open XLSX: {exc}") from exc

    warnings: list[str] = []
    sheets_data: list[dict] = []
    content_parts: list[str] = []
    truncated = False

    sheet_names = wb.sheetnames[: lim.max_sheets]
    if len(wb.sheetnames) > lim.max_sheets:
        truncated = True
        warnings.append(f"workbook has {len(wb.sheetnames)} sheets; capped at {lim.max_sheets}")

    for sheet_name in sheet_names:
        ws = wb[sheet_name]
        rows_iter = ws.iter_rows(values_only=True)
        header_row = next(rows_iter, None)
        if header_row is None:
            sheets_data.append({"name": sheet_name, "headers": [], "rows_preview": [], "row_count": 0})
            continue

        headers = [str(h) if h is not None else "" for h in header_row]
        rows_preview: list[dict] = []
        total_rows = 0
        sheet_truncated = False

        for raw_row in rows_iter:
            total_rows += 1
            if len(rows_preview) < lim.max_rows:
                rows_preview.append({
                    h: (str(v) if v is not None else "")
                    for h, v in zip(headers, raw_row)
                })
            else:
                sheet_truncated = True

        if sheet_truncated:
            truncated = True

        sheets_data.append({
            "name": sheet_name,
            "headers": headers,
            "rows_preview": rows_preview,
            "row_count": total_rows,
            "truncated": sheet_truncated,
        })

        # Content preview
        content_parts.append(f"=== {sheet_name} ===")
        content_parts.append(",".join(headers))
        for row in rows_preview:
            content_parts.append(",".join(str(row.get(h, "")) for h in headers))
        if sheet_truncated:
            content_parts.append(f"... ({total_rows - lim.max_rows} more rows)")

    wb.close()
    content = "\n".join(content_parts)
    sections = tuple(s["name"] for s in sheets_data)

    return ReadResult(
        path=path,
        format="xlsx",
        media_type=(
            "application/vnd.openxmlformats-officedocument"
            ".spreadsheetml.sheet"
        ),
        content=content,
        sections=sections,
        metadata={
            "sheet_count": len(sheets_data),
            "sheet_names": list(sections),
        },
        structured={"sheets": sheets_data},
        warnings=tuple(warnings),
        truncated=truncated,
        source_sha256=sha256,
        source_bytes=source_bytes,
    )


def read_pptx(path: str | Path, limits: ReadLimits | None = None) -> ReadResult:
    lim = limits or ReadLimits()
    path = str(path)
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError as exc:
        raise FileFormatError("python-pptx is required for PPTX reading") from exc

    try:
        raw = _read_bytes(path, lim.max_bytes)
    except OSError as exc:
        raise FileReadError(str(exc)) from exc

    sha256 = _sha256(raw)
    source_bytes = len(raw)

    try:
        prs = Presentation(io.BytesIO(raw))
    except Exception as exc:
        raise FileReadError(f"Could not open PPTX: {exc}") from exc

    warnings: list[str] = []
    slides_data: list[dict] = []
    content_parts: list[str] = []
    truncated = False
    total_slides = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        if i >= lim.max_slides:
            truncated = True
            break

        title_text = ""
        body_parts: list[str] = []
        notes_text = ""

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.shape_type == 13:  # picture
                continue
            is_title = (
                hasattr(shape, "placeholder_format")
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0
            )
            text = "\n".join(
                p.text for p in shape.text_frame.paragraphs if p.text.strip()
            )
            if is_title:
                title_text = text
            elif text:
                body_parts.append(text)

        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip() if notes_frame else ""

        body_text = "\n".join(body_parts)
        slides_data.append({
            "index": i + 1,
            "title": title_text,
            "body": body_text,
            "notes": notes_text,
        })

        slide_preview = f"--- Slide {i + 1}"
        if title_text:
            slide_preview += f": {title_text}"
        slide_preview += " ---"
        content_parts.append(slide_preview)
        if body_text:
            content_parts.append(body_text)

    content = "\n".join(content_parts)
    sections = tuple(
        s["title"] or f"Slide {s['index']}" for s in slides_data
    )

    return ReadResult(
        path=path,
        format="pptx",
        media_type=(
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation"
        ),
        content=content,
        sections=sections,
        metadata={
            "slide_count": total_slides,
            "slides_read": len(slides_data),
        },
        structured={"slides": slides_data},
        warnings=tuple(warnings),
        truncated=truncated,
        source_sha256=sha256,
        source_bytes=source_bytes,
    )


def read_md(path: str | Path, limits: ReadLimits | None = None) -> ReadResult:
    lim = limits or ReadLimits()
    path = str(path)
    try:
        raw = _read_bytes(path, lim.max_bytes)
    except OSError as exc:
        raise FileReadError(str(exc)) from exc

    sha256 = _sha256(raw)
    source_bytes = len(raw)

    warnings: list[str] = []
    try:
        text = raw.decode("utf-8-sig")
    except UnicodeDecodeError:
        text = raw.decode("latin-1", errors="replace")
        warnings.append("non-UTF-8 characters replaced")

    text = text.replace("\r\n", "\n").replace("\r", "\n")

    truncated = False
    if len(text) > lim.max_chars:
        text = text[: lim.max_chars]
        truncated = True
        warnings.append(f"content truncated at {lim.max_chars} characters")

    # Extract headings and sections
    headings: list[dict] = []
    sections_data: list[dict] = []
    current_heading: str | None = None
    current_body: list[str] = []

    import re
    heading_re = re.compile(r"^(#{1,6})\s+(.+)$")

    for line in text.splitlines():
        m = heading_re.match(line)
        if m:
            if current_heading is not None or current_body:
                sections_data.append({
                    "heading": current_heading or "",
                    "level": 0 if current_heading is None else len(m.group(1)) - 1,
                    "body": "\n".join(current_body).strip(),
                })
            current_heading = m.group(2).strip()
            headings.append({"text": current_heading, "level": len(m.group(1))})
            current_body = []
        else:
            current_body.append(line)

    if current_heading is not None or current_body:
        sections_data.append({
            "heading": current_heading or "",
            "level": headings[-1]["level"] if headings else 0,
            "body": "\n".join(current_body).strip(),
        })

    return ReadResult(
        path=path,
        format="md",
        media_type="text/markdown",
        content=text,
        sections=tuple(h["text"] for h in headings),
        metadata={
            "line_count": text.count("\n") + 1 if text else 0,
            "heading_count": len(headings),
            "char_count": len(text),
        },
        structured={"headings": headings, "sections": sections_data},
        warnings=tuple(warnings),
        truncated=truncated,
        source_sha256=sha256,
        source_bytes=source_bytes,
    )


def read_html(path: str | Path, limits: ReadLimits | None = None) -> ReadResult:
    lim = limits or ReadLimits()
    path = str(path)
    try:
        raw = _read_bytes(path, lim.max_bytes)
    except OSError as exc:
        raise FileReadError(str(exc)) from exc

    sha256 = _sha256(raw)
    source_bytes = len(raw)

    warnings: list[str] = []
    try:
        text_raw = raw.decode("utf-8-sig")
    except UnicodeDecodeError:
        text_raw = raw.decode("latin-1", errors="replace")
        warnings.append("non-UTF-8 characters replaced")

    import html as html_module
    from html.parser import HTMLParser

    class _TextExtractor(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self._skip = False
            self._skip_tags = {"script", "style", "head"}
            self.text_parts: list[str] = []
            self.headings: list[dict] = []
            self.links: list[dict] = []
            self._current_tag = ""
            self._current_text: list[str] = []

        def handle_starttag(self, tag: str, attrs: list) -> None:
            self._current_tag = tag
            if tag in self._skip_tags:
                self._skip = True
            if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                self._current_text = []
            if tag == "a":
                href = dict(attrs).get("href", "")
                if href:
                    self.links.append({"href": href, "text": ""})

        def handle_endtag(self, tag: str) -> None:
            if tag in self._skip_tags:
                self._skip = False
            if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                heading_text = "".join(self._current_text).strip()
                if heading_text:
                    self.headings.append({"tag": tag, "text": heading_text})
                    self.text_parts.append(heading_text)
                self._current_text = []

        def handle_data(self, data: str) -> None:
            if self._skip:
                return
            stripped = data.strip()
            if stripped:
                self.text_parts.append(stripped)
                if self._current_tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                    self._current_text.append(stripped)
                if self.links and not self.links[-1]["text"]:
                    self.links[-1]["text"] = stripped

    parser = _TextExtractor()
    try:
        parser.feed(text_raw)
    except Exception:
        warnings.append("HTML parse errors encountered; content may be incomplete")

    text = "\n".join(parser.text_parts)
    truncated = False
    if len(text) > lim.max_chars:
        text = text[: lim.max_chars]
        truncated = True
        warnings.append(f"extracted text truncated at {lim.max_chars} characters")

    sections = tuple(h["text"] for h in parser.headings)

    return ReadResult(
        path=path,
        format="html",
        media_type="text/html",
        content=text,
        sections=sections,
        metadata={
            "heading_count": len(parser.headings),
            "link_count": len(parser.links),
            "char_count": len(text),
        },
        structured={"headings": parser.headings, "links": parser.links},
        warnings=tuple(warnings),
        truncated=truncated,
        source_sha256=sha256,
        source_bytes=source_bytes,
    )


# ---- helpers ----

def _read_bytes(path: str, max_bytes: int) -> bytes:
    size = os.path.getsize(path)
    if size > max_bytes:
        raise FileLimitExceededError(
            f"File size {size} exceeds limit {max_bytes}"
        )
    with open(path, "rb") as f:
        return f.read()


def _sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


# Import here to avoid circular at top of module
from core.file_formats.models import FileLimitExceededError  # noqa: E402
