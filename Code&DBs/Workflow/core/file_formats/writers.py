"""File format writers: CSV, TXT, DOCX, PDF, XLSX, PPTX, Markdown, HTML."""

from __future__ import annotations

import csv
import hashlib
import os
import tempfile
from pathlib import Path

from core.file_formats.models import (
    FileFormatError,
    FileWriteError,
    FORMAT_HANDLERS,
    WriteResult,
    WriteValidationError,
)


def write_csv(
    rows: list[dict[str, object]],
    headers: list[str],
    path: str | Path,
) -> WriteResult:
    path = str(path)
    _validate_extension(path, ".csv")

    if not headers:
        raise WriteValidationError("headers must not be empty")
    if len(set(headers)) != len(headers):
        raise WriteValidationError("duplicate headers are not allowed")

    buf = _write_csv_to_bytes(rows, headers)
    return _atomic_write(path, buf, "csv")


def write_txt(content: str, path: str | Path) -> WriteResult:
    path = str(path)
    _validate_extension(path, ".txt")
    normalized = content.replace("\r\n", "\n").replace("\r", "\n")
    buf = normalized.encode("utf-8")
    return _atomic_write(path, buf, "txt")


def write_docx(
    sections: list[dict[str, str]],
    path: str | Path,
) -> WriteResult:
    path = str(path)
    _validate_extension(path, ".docx")
    try:
        import docx as python_docx
    except ImportError as exc:
        raise FileFormatError("python-docx is required for DOCX writing") from exc

    doc = python_docx.Document()
    for section in sections:
        heading = section.get("heading", "")
        body = section.get("body", "")
        if heading:
            doc.add_heading(heading, level=1)
        if body:
            for para in body.split("\n\n"):
                para = para.strip()
                if para:
                    doc.add_paragraph(para)

    buf = _docx_to_bytes(doc)
    return _atomic_write(path, buf, "docx")


def write_pdf(
    sections: list[dict[str, str]],
    path: str | Path,
) -> WriteResult:
    path = str(path)
    _validate_extension(path, ".pdf")
    try:
        from fpdf import FPDF
    except ImportError as exc:
        raise FileFormatError("fpdf2 is required for PDF writing") from exc

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)

    for section in sections:
        heading = section.get("heading", "")
        body = section.get("body", "")
        if heading:
            pdf.set_font("Helvetica", style="B", size=14)
            # fpdf2 replaces unmappable chars; use errors-tolerant encode
            safe_heading = _sanitize_text(heading)
            pdf.multi_cell(0, 8, safe_heading)
            pdf.ln(2)
            pdf.set_font("Helvetica", size=12)
        if body:
            for para in body.split("\n\n"):
                para = para.strip()
                if para:
                    safe_para = _sanitize_text(para)
                    pdf.multi_cell(0, 6, safe_para)
                    pdf.ln(4)

    buf = pdf.output()
    if isinstance(buf, str):
        buf = buf.encode("latin-1")
    return _atomic_write(path, bytes(buf), "pdf")


def write_xlsx(
    sheets: list[dict],
    path: str | Path,
) -> WriteResult:
    """Write one or more sheets to an XLSX file.

    Each sheet dict: {"name": str, "headers": list[str], "rows": list[dict]}
    """
    path = str(path)
    _validate_extension(path, ".xlsx")
    try:
        import openpyxl
        from openpyxl.styles import Font
    except ImportError as exc:
        raise FileFormatError("openpyxl is required for XLSX writing") from exc

    if not sheets:
        raise WriteValidationError("sheets must not be empty")

    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # remove default sheet

    for sheet in sheets:
        name = sheet.get("name", "Sheet")
        headers = sheet.get("headers", [])
        rows = sheet.get("rows", [])

        if len(set(headers)) != len(headers):
            raise WriteValidationError(f"duplicate headers in sheet {name!r}")

        ws = wb.create_sheet(title=name)
        if headers:
            ws.append(headers)
            for cell in ws[1]:
                cell.font = Font(bold=True)
        for row in rows:
            ws.append([row.get(h, "") for h in headers])

    buf = _wb_to_bytes(wb)
    return _atomic_write(path, buf, "xlsx")


def write_pptx(
    slides: list[dict[str, str]],
    path: str | Path,
) -> WriteResult:
    """Write slides to a PPTX file.

    Each slide dict: {"title": str, "body": str}  (both optional)
    """
    path = str(path)
    _validate_extension(path, ".pptx")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise FileFormatError("python-pptx is required for PPTX writing") from exc

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]       # blank
    title_content_layout = prs.slide_layouts[1]  # title + content

    for slide_data in slides:
        title = slide_data.get("title", "")
        body = slide_data.get("body", "")

        if title and body:
            slide = prs.slides.add_slide(title_content_layout)
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, para in enumerate(body.split("\n\n")):
                para = para.strip()
                if not para:
                    continue
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = para
        elif title:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
        elif body:
            slide = prs.slides.add_slide(blank_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            txBox.text_frame.text = body

    import io as _io
    buf = _io.BytesIO()
    prs.save(buf)
    return _atomic_write(path, buf.getvalue(), "pptx")


def write_md(
    sections: list[dict[str, str]],
    path: str | Path,
) -> WriteResult:
    """Write sections to a Markdown file.

    Each section dict: {"heading": str, "level": int (1-6), "body": str}
    heading and level are optional.
    """
    path = str(path)
    _validate_extension(path, ".md")
    parts: list[str] = []
    for section in sections:
        heading = section.get("heading", "")
        level = max(1, min(6, int(section.get("level", 1))))
        body = section.get("body", "")
        if heading:
            parts.append(f"{'#' * level} {heading}\n")
        if body:
            parts.append(body.strip())
            parts.append("")
    content = "\n".join(parts).rstrip() + "\n"
    buf = content.encode("utf-8")
    return _atomic_write(path, buf, "md")


def write_html(
    sections: list[dict[str, str]],
    path: str | Path,
    title: str = "",
) -> WriteResult:
    """Write sections to an HTML file.

    Each section dict: {"heading": str, "level": int (1-6), "body": str}
    """
    path = str(path)
    ext = os.path.splitext(path)[1].lower()
    if ext not in (".html", ".htm"):
        raise WriteValidationError(
            f"path extension {ext!r} does not match expected .html or .htm"
        )
    import html as _html
    parts: list[str] = [
        "<!DOCTYPE html>",
        "<html>",
        "<head>",
        f"<meta charset='utf-8'>",
        f"<title>{_html.escape(title or 'Document')}</title>",
        "</head>",
        "<body>",
    ]
    for section in sections:
        heading = section.get("heading", "")
        level = max(1, min(6, int(section.get("level", 2))))
        body = section.get("body", "")
        if heading:
            parts.append(f"<h{level}>{_html.escape(heading)}</h{level}>")
        if body:
            for para in body.split("\n\n"):
                para = para.strip()
                if para:
                    parts.append(f"<p>{_html.escape(para)}</p>")
    parts += ["</body>", "</html>"]
    content = "\n".join(parts) + "\n"
    buf = content.encode("utf-8")
    return _atomic_write(path, buf, "html")


# ---- helpers ----

def _validate_extension(path: str, expected: str) -> None:
    ext = os.path.splitext(path)[1].lower()
    if ext != expected:
        raise WriteValidationError(
            f"path extension {ext!r} does not match expected {expected!r}"
        )


def _atomic_write(path: str, data: bytes, fmt: str) -> WriteResult:
    ext = os.path.splitext(path)[1].lower()
    handler = FORMAT_HANDLERS.get(ext) or FORMAT_HANDLERS.get(f".{fmt}")
    media_type = handler.media_type if handler else "application/octet-stream"
    tmp = path + ".tmp"
    try:
        with open(tmp, "wb") as f:
            f.write(data)
        os.replace(tmp, path)
    except OSError as exc:
        # Clean up tmp if it exists
        try:
            os.unlink(tmp)
        except OSError:
            pass
        raise FileWriteError(str(exc)) from exc

    sha = hashlib.sha256(data).hexdigest()
    return WriteResult(
        path=path,
        format=fmt,  # type: ignore[arg-type]
        media_type=media_type,
        bytes_written=len(data),
        sha256=sha,
        metadata={},
    )


def _write_csv_to_bytes(rows: list[dict[str, object]], headers: list[str]) -> bytes:
    buf = []
    sio = _StringIOCollector(buf)
    writer = csv.DictWriter(
        sio,
        fieldnames=headers,
        lineterminator="\n",
        extrasaction="raise",
        restval="",
    )
    writer.writeheader()
    for row in rows:
        writer.writerow(row)
    return "".join(buf).encode("utf-8")


def _wb_to_bytes(wb: object) -> bytes:
    import io as _io
    buf = _io.BytesIO()
    wb.save(buf)  # type: ignore[attr-defined]
    return buf.getvalue()


def _docx_to_bytes(doc: object) -> bytes:
    import io as _io
    buf = _io.BytesIO()
    doc.save(buf)  # type: ignore[attr-defined]
    return buf.getvalue()


def _sanitize_text(text: str) -> str:
    # Replace characters outside latin-1 range with '?' for fpdf2 core fonts
    return text.encode("latin-1", errors="replace").decode("latin-1")


class _StringIOCollector:
    """Minimal write-only file-like object that collects to a list."""

    def __init__(self, buf: list[str]) -> None:
        self._buf = buf

    def write(self, s: str) -> int:
        self._buf.append(s)
        return len(s)
